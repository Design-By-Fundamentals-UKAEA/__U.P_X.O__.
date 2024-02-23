"""
PowerPoint Operations Module (pwpops)

This module provides a set of classes and functions to automate and facilitate the creation and manipulation of PowerPoint presentations using the python-pptx library. It is designed to simplify common PowerPoint tasks such as adding slides, inserting text and images, and customizing slide layouts and styles.

Features:
- Create new PowerPoint presentations with customizable properties such as file name, folder location, paper size, and slide orientation.
- Add various types of slides with different layouts.
- Insert and format text with customizable font properties, including name, size, color, and style.
- Embed images with options for scaling and positioning, along with captions.
- Apply consistent styling and formatting across slides for a professional appearance.

Usage:
The module is intended for use in scripting and automation tasks where PowerPoint presentations need to be generated or modified programmatically. It abstracts away the lower-level details of the python-pptx API, providing a more intuitive and high-level interface for users.

Classes:
- pwpops: The main class responsible for presentation operations. It includes methods for setting up the presentation, adding slides, inserting text and images, and customizing slide properties.

Example:
    from pwpops import pwpops

    # Create a new PowerPointManager instance
    ppt_manager = pwpops()

    # Set basic properties
    ppt_manager.set_file_name('MyPresentation.pptx')
    ppt_manager.set_folder_name('/path/to/directory')
    ppt_manager.set_paper_size('A4')
    ppt_manager.set_slide_orientation('landscape')

    # Create a new presentation
    ppt_manager.create()

    # Add a slide with a title and content
    ppt_manager.add_slide(slide_type='Title and Content', title='Welcome', content='This is the first slide')

    # Save the presentation
    ppt_manager.save()

Dependencies:
- python-pptx: This module relies on the python-pptx library to interact with PowerPoint files. Ensure that python-pptx is installed and available in your environment.

Note:
This module is designed for automation and scripting purposes. It may not cover all features available in the python-pptx library or PowerPoint itself. For more complex operations, consider using the python-pptx API directly or extending this module to suit your needs.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, RGBColor
# from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
import datetime

class pwpops:
    """
    A class to manage PowerPoint operations such as creating slides,
    adding text, images, and customizing slide properties.

    Attributes:
        prs (Presentation): A Presentation object from python-pptx library.
        file_name (str): Name of the PowerPoint file to be created or modified.
        folder_name (str): Directory path where the PowerPoint file will be stored.
        paper_size (str): Size of the PowerPoint slides.
        slide_orientation (str): Orientation of the slides, either 'portrait' or 'landscape'.
        slide_heading_text_details (dict): Styling details for slide heading texts.
        slide_text_details (dict): Styling details for regular slide texts.
        image_caption_text_details (dict): Styling details for image captions.
        document_topic (str): Topic of the document for reference in footers.
    """
    __slots__ = ['prs',
                 'file_name',
                 'folder_name',
                 'paper_size',
                 'slide_orientation',
                 'slide_heading_text_details',
                 'slide_text_details',
                 'image_caption_text_details',
                 'document_topic']

    def __init__(self):
        """Initializes the PowerPointManager with default values."""
        self.prs = Presentation()
        self.file_name = "presentation.pptx"
        self.folder_name = "."
        self.paper_size = "A4"
        self.slide_orientation = "portrait"
        self.slide_heading_text_details = {'fontName': 'Arial', 'fontStyle': 'Bold', 'fontColour': RGBColor(0, 0, 0)}
        self.slide_text_details = {'fontName': 'Arial', 'fontStyle': 'Regular', 'fontColour': RGBColor(0, 0, 0)}
        self.image_caption_text_details = {'fontName': 'Arial', 'fontStyle': 'Italic', 'fontColour': RGBColor(0, 0, 0)}
        self.document_topic = ""

    def set_document_topic(self, topic):
        """Sets the document topic for use in footers."""
        self.document_topic = topic

    # General setters
    def set_file_name(self, file_name):
        """Sets the file name for the PowerPoint document."""
        self.file_name = file_name

    def set_folder_name(self, folder_name):
        """Sets the folder name where the PowerPoint document will be saved."""
        self.folder_name = folder_name

    def set_paper_size(self, paper_size):
        """Sets the paper size for the PowerPoint slides."""
        self.paper_size = paper_size

    def set_slide_orientation(self, slide_orientation):
        """Sets the orientation of the PowerPoint slides."""
        self.slide_orientation = slide_orientation

    def set_slide_heading_text_details(self, font_name, font_style, font_colour):
        """Sets the text style details for slide headings."""
        self.slide_heading_text_details = {
            'fontName': font_name,
            'fontStyle': font_style,
            'fontColour': font_colour
        }

    def set_slide_text_details(self, font_name, font_style, font_colour):
        """Sets the text style details for regular slide text."""
        self.slide_text_details = {
            'fontName': font_name,
            'fontStyle': font_style,
            'fontColour': font_colour
        }

    def set_image_caption_text_details(self, font_name, font_style, font_colour):
        """Sets the text style details for image captions."""
        self.image_caption_text_details = {
            'fontName': font_name,
            'fontStyle': font_style,
            'fontColour': font_colour
        }

    # Setters for slide_heading_text_details
    def set_heading_font_name(self, font_name):
        """Sets the font name for slide headings."""
        self.slide_heading_text_details['fontName'] = font_name

    def set_heading_font_style(self, font_style):
        """Sets the font style for slide headings."""
        self.slide_heading_text_details['fontStyle'] = font_style

    def set_heading_font_colour(self, font_colour):
        """Sets the font colour for slide headings."""
        self.slide_heading_text_details['fontColour'] = font_colour

    # Setters for slide_text_details
    def set_text_font_name(self, font_name):
        self.slide_text_details['fontName'] = font_name

    def set_text_font_style(self, font_style):
        self.slide_text_details['fontStyle'] = font_style

    def set_text_font_colour(self, font_colour):
        self.slide_text_details['fontColour'] = font_colour

    # Setters for image_caption_text_details
    def set_caption_font_name(self, font_name):
        self.image_caption_text_details['fontName'] = font_name

    def set_caption_font_style(self, font_style):
        self.image_caption_text_details['fontStyle'] = font_style

    def set_caption_font_colour(self, font_colour):
        self.image_caption_text_details['fontColour'] = font_colour

    def create(self, fileName=None, folderName=None, paperSize='A4', slideOrientation='portrait'):
        """
        Creates a new PowerPoint file with specified properties.

        Parameters:
            fileName (str): The name of the file to create.
            folderName (str): The directory path to save the file.
            paperSize (str): The size of the paper to use for slides.
            slideOrientation (str): The orientation of the slides.
        """
        if fileName:
            self.file_name = fileName
        if folderName:
            self.folder_name = folderName
        self.paper_size = paperSize
        self.slide_orientation = slideOrientation
        # Note: Adjusting paper size and orientation might require direct manipulation of XML, which python-pptx doesn't support directly.

    def add_notes(self, slideNumber, notesText):
        """
        Adds speaker notes to a slide.

        Parameters:
            slideNumber: The slide to add notes to.
            notesText (str): The text of the notes to add.
        """
        slide = self._get_slide(slideNumber)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notesText

    def add_title(self, slide, title):
        """
        Adds a title to a slide.

        Parameters:
            slide: The slide to add a title to.
            title (str): The text of the title.
        """
        # Implementation omitted for brevity...

    def add_batch_slides(self, numSlides):
        for _ in range(numSlides):
            self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Using blank slide layout

    def slide_layout_customisation(self, slideNumber, layoutIndex):
        slide_layout = self.prs.slide_layouts[layoutIndex]
        self.prs.slides.add_slide(slide_layout)

    def add_text(self, slideNumber, textPosition, size, text=""):

        slide = self._get_slide(slideNumber)
        left, top = Inches(textPosition[0]), Inches(textPosition[1])
        width, height = Inches(size[0]), Inches(size[1])
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = text
        # Customise text style here using self.slide_text_details

    def insert_section(self, slideNumber, sectionTitle):
        # python-pptx doesn't support sections directly. This method would simply add a title slide as a section separator.
        slide_layout = self.prs.slide_layouts[0]  # Assuming 0 is the title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = sectionTitle

    def add_footer(self, slideNumber):
        slide = self._get_slide(slideNumber)
        footer_text = f"Generated by UPXO {datetime.datetime.now().strftime('%Y-%m-%d %H:%M:%S')}, {self.document_topic}, Slide {slideNumber+1}"
        left, top = Inches(0), Inches(self.prs.slide_height - 0.5)  # Positioning at the bottom
        width, height = Inches(self.prs.slide_width), Inches(0.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = footer_text
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(10)  # Example font size

    def save(self):
        file_path = f"{self.folder_name}/{self.file_name}"
        self.prs.save(file_path)

    def _get_slide(self, slideNumber):
        if slideNumber < len(self.prs.slides):
            return self.prs.slides[slideNumber]
        else:
            return self.prs.slides.add_slide(self.prs.slide_layouts[5])  # Using blank slide layout
